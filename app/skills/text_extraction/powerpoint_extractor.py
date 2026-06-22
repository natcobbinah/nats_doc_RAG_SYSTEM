import io 
import logging
import os 
import sys 
from I_document_handler import IDocumentHandler
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict
from pptx import Presentation

if __package__:
    from ...logging_utils import configure_json_logging, log, LogLevel
    from ...utils import sanitize_and_standardize_doc_id
else:
    sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..")))
    from app.logging_utils import configure_json_logging, log, LogLevel
    from app.utils import sanitize_and_standardize_doc_id

class TextFromPowerPoint(IDocumentHandler):
    def __init__(self):
        configure_json_logging()

    def _process_slide(self, args):
        slide, i, blob_name = args
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        return i, text
    
    def extract_content_from_document(self, stream: io.BytesIO, blob_name: str, source_url: str) -> List[Dict]:
        """Extract text and metadata from PowerPoint, one page per slide"""
        try:
            presentation = Presentation(stream)
            props = presentation.core_properties
            total_slides = len(presentation.slides)

            doc_metadata = {
                "source_filename": blob_name,
                "title": props.title or "",
                "author": props.author or "",
                "creation_date": props.created.isoformat() if props.created else "",
                "modification_date": props.modified.isoformat() if props.modified else "",
                "total_pages": total_slides,
                "file_extension": "pptx",
                "url": source_url
            }

            slides_data = []

            with ThreadPoolExecutor() as executor:
                futures = [executor.submit(self._process_slide, (presentation.slides[i],i, blob_name)) for i in range(total_slides)]
                for future in as_completed(futures):
                    i, text = future.result()
                    slides_data.append({
                        "id": sanitize_and_standardize_doc_id(f"{blob_name}-slide-{int(i + 1)}"),
                        "page_number": int(i + 1),
                        "content": text, 
                        **doc_metadata
                    })
            
            # sort slides by page number
            slides_data = sorted(slides_data, key=lambda x: int(x["page_number"]))

            return slides_data
        
        except Exception as e:
            log(
                type=LogLevel.ERROR,
                message = {
                    "message":f"Failed to extract  content from  word document, {blob_name}",
                    "fileline": sys._getframe().f_lineno,
                    "message_data": {"error": str(e)},
                    "reason": str(e)
                }
            )
            raise

if __name__ == "__main__":
    dummy_file_path = os.path.join(os.path.dirname(__file__), "diapo.pptx")
    with open(dummy_file_path, "rb") as file_stream:
        tft = TextFromPowerPoint()
        result = tft.extract_content_from_document(
            stream=io.BytesIO(file_stream.read()),
            blob_name="diapo",
            source_url=""
        )
        print(result)

 